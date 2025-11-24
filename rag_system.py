#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Sistema RAG (Retrieval-Augmented Generation) para Chatbot CAMACOL
Procesa documentos PDF, Word, Excel, CSV, PowerPoint y realiza búsqueda semántica con prioridad por año
"""

import os
import pickle
import requests
import tempfile
from pathlib import Path
from typing import List, Tuple, Optional, Dict, Any
import warnings
warnings.filterwarnings('ignore')

# Procesamiento de documentos
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except:
    PPTX_AVAILABLE = False

# Vector store
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_community.vectorstores import FAISS
    from langchain_community.embeddings import HuggingFaceEmbeddings
    from langchain_core.documents import Document as LangchainDocument
    VECTORSTORE_AVAILABLE = True
except Exception as e:
    VECTORSTORE_AVAILABLE = False
    print(f"⚠️ Vector store no disponible: {e}")

class RAGSystem:
    """Sistema RAG para procesamiento y búsqueda de documentos con prioridad por año"""
    
    def __init__(self, rag_folder: str, cache_folder: str = "./rag_cache"):
        self.rag_folder = Path(rag_folder)
        self.cache_folder = Path(cache_folder)
        self.cache_folder.mkdir(exist_ok=True)
        
        self.vectorstore = None
        self.embeddings = None
        self.documents = []
        self.metadata = []
        
        # Prioridad de años (más reciente primero)
        self.year_priority = ['2025', '2024', '2023']
        
        self.available = VECTORSTORE_AVAILABLE and (PDF_AVAILABLE or DOCX_AVAILABLE or PANDAS_AVAILABLE or PPTX_AVAILABLE)
    
    def _descargar_desde_url(self, url: str, carpeta_destino: Path) -> Optional[Path]:
        """Descarga un archivo desde una URL"""
        try:
            print(f"🌐 Descargando: {url}")
            response = requests.get(url, timeout=30, allow_redirects=True)
            response.raise_for_status()
            
            # Obtener nombre del archivo desde la URL o Content-Disposition
            filename = url.split('/')[-1]
            if 'Content-Disposition' in response.headers:
                content_disp = response.headers['Content-Disposition']
                if 'filename=' in content_disp:
                    filename = content_disp.split('filename=')[1].strip('"')
            
            # Guardar archivo
            archivo_path = carpeta_destino / filename
            with open(archivo_path, 'wb') as f:
                f.write(response.content)
            
            print(f"✅ Descargado: {filename} ({len(response.content)} bytes)")
            return archivo_path
            
        except Exception as e:
            print(f"❌ Error descargando {url}: {e}")
            return None
    
    def _procesar_urls_desde_archivo(self, archivo_urls: Path) -> int:
        """Procesa URLs desde un archivo de texto"""
        if not archivo_urls.exists():
            print(f"⚠️ Archivo de URLs no encontrado: {archivo_urls}")
            return 0
        
        archivos_procesados = 0
        
        # Crear carpeta temporal para descargas
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # Leer URLs del archivo
            with open(archivo_urls, 'r', encoding='utf-8') as f:
                urls = [line.strip() for line in f if line.strip() and not line.startswith('#')]
            
            print(f"📚 Procesando {len(urls)} URLs...")
            
            # Descargar y procesar cada URL
            for url in urls:
                archivo_descargado = self._descargar_desde_url(url, temp_path)
                
                if archivo_descargado and archivo_descargado.exists():
                    try:
                        ext = archivo_descargado.suffix.lower()
                        
                        if ext == '.pdf' and PDF_AVAILABLE:
                            self._procesar_pdf(archivo_descargado)
                            archivos_procesados += 1
                        elif ext in ['.docx', '.doc'] and DOCX_AVAILABLE:
                            self._procesar_word(archivo_descargado)
                            archivos_procesados += 1
                        elif ext in ['.xlsx', '.xls'] and PANDAS_AVAILABLE:
                            self._procesar_excel(archivo_descargado)
                            archivos_procesados += 1
                        elif ext == '.csv' and PANDAS_AVAILABLE:
                            self._procesar_csv(archivo_descargado)
                            archivos_procesados += 1
                        elif ext in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
                            self._procesar_pptx(archivo_descargado)
                            archivos_procesados += 1
                        else:
                            print(f"⚠️ Formato no soportado: {ext}")
                    except Exception as e:
                        print(f"❌ Error procesando {archivo_descargado.name}: {e}")
        
        return archivos_procesados
    
    def inicializar(self, force_reload: bool = False) -> Tuple[bool, str]:
        """Inicializa el sistema RAG"""
        if not self.available:
            return False, "❌ Sistema RAG no disponible"
        
        try:
            cache_file = self.cache_folder / "vectorstore.pkl"
            
            if cache_file.exists() and not force_reload:
                print("💾 Cargando vector store desde cache...")
                with open(cache_file, 'rb') as f:
                    cache_data = pickle.load(f)
                    self.vectorstore = cache_data['vectorstore']
                    self.metadata = cache_data['metadata']
                
                return True, f"✅ Cache cargado: {len(self.metadata)} documentos"
            
            # Procesar documentos
            print("📂 Procesando documentos locales...")
            exito, mensaje = self._procesar_documentos()
            
            if not exito:
                return False, mensaje
            
            # Procesar URLs si existe el archivo
            archivo_urls = Path("urls_documentos.txt")
            if archivo_urls.exists():
                print("🌐 Procesando documentos desde URLs...")
                archivos_url = self._procesar_urls_desde_archivo(archivo_urls)
                if archivos_url > 0:
                    print(f"✅ Procesados {archivos_url} archivos desde URLs")
            
            # Crear embeddings
            print("🧠 Creando embeddings...")
            self.embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
            )
            
            # Crear vector store
            print("📊 Creando vector store...")
            if len(self.documents) == 0:
                return False, "❌ No se encontraron documentos"
            
            self.vectorstore = FAISS.from_documents(self.documents, self.embeddings)
            
            # Guardar cache
            with open(cache_file, 'wb') as f:
                pickle.dump({
                    'vectorstore': self.vectorstore,
                    'metadata': self.metadata
                }, f)
            
            return True, f"✅ RAG inicializado: {len(self.metadata)} docs, {len(self.documents)} chunks"
            
        except Exception as e:
            return False, f"❌ Error: {str(e)}"
    
    def _procesar_documentos(self) -> Tuple[bool, str]:
        """Procesa todos los documentos con prioridad por año"""
        if not self.rag_folder.exists():
            return False, f"❌ Carpeta no encontrada: {self.rag_folder}"
        
        archivos_procesados = 0
        
        # Procesar por año en orden de prioridad
        for year in self.year_priority:
            year_folder = self.rag_folder / year
            if year_folder.exists():
                print(f"📅 Procesando año: {year}")
                archivos_procesados += self._procesar_carpeta(year_folder)
        
        # Procesar archivos en la raíz (si existen)
        archivos_procesados += self._procesar_carpeta(self.rag_folder, solo_raiz=True)
        
        if archivos_procesados == 0:
            return False, "❌ No se procesaron archivos"
        
        return True, f"✅ Procesados {archivos_procesados} archivos"
    
    def _procesar_carpeta(self, carpeta: Path, solo_raiz: bool = False) -> int:
        """Procesa archivos en una carpeta"""
        archivos_procesados = 0
        
        patron = '*' if solo_raiz else '**/*'
        for archivo in carpeta.glob(patron):
            if archivo.is_file():
                try:
                    ext = archivo.suffix.lower()
                    
                    if ext == '.pdf' and PDF_AVAILABLE:
                        self._procesar_pdf(archivo)
                        archivos_procesados += 1
                    elif ext in ['.docx', '.doc'] and DOCX_AVAILABLE:
                        self._procesar_word(archivo)
                        archivos_procesados += 1
                    elif ext in ['.xlsx', '.xls'] and PANDAS_AVAILABLE:
                        self._procesar_excel(archivo)
                        archivos_procesados += 1
                    elif ext == '.csv' and PANDAS_AVAILABLE:
                        self._procesar_csv(archivo)
                        archivos_procesados += 1
                    elif ext in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
                        self._procesar_pptx(archivo)
                        archivos_procesados += 1
                except Exception as e:
                    print(f"❌ Error en {archivo.name}: {e}")
        
        return archivos_procesados
    
    def _procesar_pdf(self, archivo: Path):
        """Procesa PDF"""
        reader = PdfReader(str(archivo))
        texto = "\n".join([p.extract_text() for p in reader.pages if p.extract_text()])
        
        if texto.strip():
            self._agregar_documento(texto, archivo, "PDF")
            print(f"✅ PDF: {archivo.name}")
    
    def _procesar_word(self, archivo: Path):
        """Procesa Word"""
        doc = Document(str(archivo))
        texto = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        
        if texto.strip():
            self._agregar_documento(texto, archivo, "Word")
            print(f"✅ Word: {archivo.name}")
    
    def _procesar_excel(self, archivo: Path):
        """Procesa Excel - para análisis de datos usar Data Analyzer"""
        excel_file = pd.ExcelFile(str(archivo))
        texto = f"Excel: {archivo.name}\n\n"
        
        for sheet in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet)
            texto += f"Hoja: {sheet}\nColumnas: {', '.join(df.columns)}\nFilas: {len(df)}\n"
            texto += df.head(10).to_string() + "\n\n"
        
        if texto.strip():
            self._agregar_documento(texto, archivo, "Excel")
            print(f"✅ Excel: {archivo.name}")
    
    def _procesar_csv(self, archivo: Path):
        """Procesa CSV - para análisis de datos usar Data Analyzer"""
        try:
            df = pd.read_csv(str(archivo))
            texto = f"CSV: {archivo.name}\n\n"
            texto += f"Columnas: {', '.join(df.columns)}\nFilas: {len(df)}\n"
            texto += df.head(10).to_string() + "\n\n"
            
            if texto.strip():
                self._agregar_documento(texto, archivo, "CSV")
                print(f"✅ CSV: {archivo.name}")
        except Exception as e:
            print(f"❌ Error procesando CSV {archivo.name}: {e}")
    
    def _procesar_pptx(self, archivo: Path):
        """Procesa PowerPoint"""
        try:
            prs = Presentation(str(archivo))
            texto = f"PowerPoint: {archivo.name}\n\n"
            
            for i, slide in enumerate(prs.slides, 1):
                texto += f"--- Diapositiva {i} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texto += shape.text + "\n"
                texto += "\n"
            
            if texto.strip():
                self._agregar_documento(texto, archivo, "PowerPoint")
                print(f"✅ PowerPoint: {archivo.name}")
        except Exception as e:
            print(f"❌ Error procesando PPTX {archivo.name}: {e}")
    
    def _agregar_documento(self, texto: str, archivo: Path, tipo: str):
        """Divide y agrega documento"""
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_text(texto)
        
        for i, chunk in enumerate(chunks):
            metadata = {
                "source": str(archivo),
                "filename": archivo.name,
                "type": tipo,
                "chunk": i,
                "folder": archivo.parent.name
            }
            doc = LangchainDocument(page_content=chunk, metadata=metadata)
            self.documents.append(doc)
            self.metadata.append(metadata)
    
    def buscar(self, query: str, k: int = 5) -> Tuple[bool, List[Dict]]:
        """Búsqueda semántica"""
        if not self.vectorstore:
            return False, []
        
        try:
            docs = self.vectorstore.similarity_search(query, k=k)
            resultados = [{"content": doc.page_content, "metadata": doc.metadata} for doc in docs]
            return True, resultados
        except Exception as e:
            return False, []
    
    def obtener_archivos_datos(self, year: Optional[str] = None) -> List[Path]:
        """Obtiene archivos de datos (Excel/CSV) priorizados por año, con LIVO como prioridad absoluta"""
        archivos_datos = []
        
        # PRIORIDAD ABSOLUTA: LIVO_total_oct25_.xlsx
        livo_path = self.rag_folder / "2025" / "Coordenada Urbana" / "LIVO_total_oct25_.xlsx"
        if livo_path.exists():
            archivos_datos.append(livo_path)
            print(f"✅ LIVO encontrado (prioridad absoluta): {livo_path.name}")
        
        # Si se especifica un año, buscar solo en ese año
        if year:
            year_folder = self.rag_folder / year
            if year_folder.exists():
                for archivo in year_folder.glob('**/*.xlsx'):
                    if archivo not in archivos_datos:
                        archivos_datos.append(archivo)
                for archivo in year_folder.glob('**/*.xls'):
                    if archivo not in archivos_datos:
                        archivos_datos.append(archivo)
                for archivo in year_folder.glob('**/*.csv'):
                    if archivo not in archivos_datos:
                        archivos_datos.append(archivo)
        else:
            # Buscar por prioridad de año
            for year_p in self.year_priority:
                year_folder = self.rag_folder / year_p
                if year_folder.exists():
                    for archivo in year_folder.glob('**/*.xlsx'):
                        if archivo not in archivos_datos:
                            archivos_datos.append(archivo)
                    for archivo in year_folder.glob('**/*.xls'):
                        if archivo not in archivos_datos:
                            archivos_datos.append(archivo)
                    for archivo in year_folder.glob('**/*.csv'):
                        if archivo not in archivos_datos:
                            archivos_datos.append(archivo)
            
            # Agregar archivos en la raíz
            for archivo in self.rag_folder.glob('*.xlsx'):
                if archivo not in archivos_datos:
                    archivos_datos.append(archivo)
            for archivo in self.rag_folder.glob('*.xls'):
                if archivo not in archivos_datos:
                    archivos_datos.append(archivo)
            for archivo in self.rag_folder.glob('*.csv'):
                if archivo not in archivos_datos:
                    archivos_datos.append(archivo)
        
        return archivos_datos
    
    def buscar_con_analisis(self, query: str, k: int = 5) -> Tuple[bool, Dict[str, Any]]:
        """Búsqueda híbrida: RAG + identificación de archivos para análisis"""
        resultado = {
            "rag_results": [],
            "data_files": [],
            "needs_analysis": False,
            "year_detected": None
        }
        
        # 1. Realizar búsqueda RAG
        exito, rag_results = self.buscar(query, k=k)
        if exito:
            resultado["rag_results"] = rag_results
        
        # 2. Detectar si necesita análisis de datos (LÓGICA MEJORADA)
        # Requiere OPERACIÓN + DATO para evitar falsos positivos
        
        query_lower = query.lower()
        
        # OPERACIONES que indican análisis de datos
        operaciones = [
            'suma', 'sumar', 'total', 'totales', 'totalizar',
            'cuenta', 'contar', 'cuántos', 'cuántas', 'cuánto',
            'promedio', 'media', 'calcular', 'cálculo',
            'compara', 'comparar', 'comparación',
            'agrupa', 'agrupar', 'agrupación',
            'filtra', 'filtrar', 'filtro',
            'top', 'ranking', 'mayor', 'menor', 'máximo', 'mínimo',
            'porcentaje', 'incremento', 'variación', 'crecimiento',
            'lista', 'listar', 'mostrar', 'dame'
        ]
        
        # DATOS específicos del sector constructor
        datos_sector = [
            'licencia', 'licencias', 'livo',
            'unidades', 'proyectos', 'proyecto',
            'ventas', 'venta', 'valor', 'valores',
            'área', 'áreas', 'metros', 'm2', 'cuadrados',
            'registros', 'datos', 'cifras',
            'vivienda', 'viviendas', 'edificación',
            'vis', 'vip', 'vur'
        ]
        
        # Detectar si hay OPERACIÓN + DATO
        tiene_operacion = any(op in query_lower for op in operaciones)
        tiene_dato = any(dato in query_lower for dato in datos_sector)
        
        # CASOS ESPECIALES: Menciones explícitas de LIVO o archivos de datos
        menciona_livo = 'livo' in query_lower
        menciona_datos = any(palabra in query_lower for palabra in ['tabla', 'excel', 'csv', 'base de datos', 'registros'])
        
        # DECISIÓN: Necesita análisis si:
        # 1. Tiene OPERACIÓN + DATO (ej: "suma licencias", "total metros")
        # 2. O menciona explícitamente LIVO o archivos de datos
        resultado["needs_analysis"] = (tiene_operacion and tiene_dato) or menciona_livo or menciona_datos
        
        # DEBUG: Mostrar por qué se activó/desactivó el análisis
        if resultado["needs_analysis"]:
            razones = []
            if tiene_operacion and tiene_dato:
                razones.append("Operación + Dato detectado")
            if menciona_livo:
                razones.append("Menciona LIVO")
            if menciona_datos:
                razones.append("Menciona archivos de datos")
            print(f"📊 Análisis activado: {', '.join(razones)}")
        else:
            print(f"📚 Análisis NO activado - Buscará en documentos RAG")
        
        # 3. Detectar año en la consulta
        for year in self.year_priority:
            if year in query:
                resultado["year_detected"] = year
                break
        
        # 4. Si necesita análisis, obtener archivos de datos relevantes
        if resultado["needs_analysis"]:
            # Primero intentar con año detectado
            if resultado["year_detected"]:
                resultado["data_files"] = self.obtener_archivos_datos(resultado["year_detected"])
            
            # Si no hay archivos o no se detectó año, buscar en todos priorizando
            if not resultado["data_files"]:
                resultado["data_files"] = self.obtener_archivos_datos()
            
            # Limitar a los primeros 5 archivos más relevantes
            resultado["data_files"] = resultado["data_files"][:5]
        
        return True, resultado
    
    def obtener_contexto(self, query: str, k: int = 3) -> str:
        """Obtiene contexto para RAG"""
        exito, resultados = self.buscar(query, k=k)
        
        if not exito or not resultados:
            return "No se encontró información relevante."
        
        contexto = "### 📚 Información de documentos CAMACOL:\n\n"
        
        for i, res in enumerate(resultados, 1):
            meta = res['metadata']
            contexto += f"**Doc {i}: {meta['filename']}** ({meta['type']})\n"
            contexto += f"{res['content']}\n\n---\n\n"
        
        return contexto
    
    def listar_documentos(self) -> str:
        """Lista documentos indexados"""
        if not self.metadata:
            return "❌ No hay documentos"
        
        archivos = {}
        for meta in self.metadata:
            fn = meta['filename']
            if fn not in archivos:
                archivos[fn] = {'tipo': meta['type'], 'folder': meta['folder'], 'chunks': 0}
            archivos[fn]['chunks'] += 1
        
        resultado = f"📚 **{len(archivos)} documentos indexados**\n\n"
        
        por_carpeta = {}
        for fn, info in archivos.items():
            folder = info['folder']
            if folder not in por_carpeta:
                por_carpeta[folder] = []
            por_carpeta[folder].append((fn, info))
        
        for folder, docs in por_carpeta.items():
            resultado += f"### 📁 {folder}\n"
            for fn, info in docs:
                resultado += f"- **{fn}** ({info['tipo']}) - {info['chunks']} chunks\n"
            resultado += "\n"
        
        return resultado